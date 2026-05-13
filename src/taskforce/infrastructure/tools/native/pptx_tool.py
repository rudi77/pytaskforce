"""PPTX presentation manipulation tool.

Provides reading, writing, and editing capabilities for Microsoft PowerPoint
(.pptx) files. ``python-pptx`` ships as a core dependency.
"""

from __future__ import annotations

from pathlib import Path
from typing import Any

import structlog

from taskforce.core.interfaces.tools import ApprovalRiskLevel
from taskforce.infrastructure.tools.base_tool import BaseTool

logger = structlog.get_logger(__name__)


def _require_pptx() -> Any:
    """Lazily import python-pptx and raise a clear error if missing."""
    try:
        import pptx

        return pptx
    except ImportError:
        raise ImportError(
            "python-pptx is part of the core install but the import failed. "
            "Run 'uv sync' to repair the venv."
        )


class PptxTool(BaseTool):
    """Read, create, and edit Microsoft PowerPoint (.pptx) presentations."""

    tool_name = "pptx"
    tool_description = (
        "Manipulate Microsoft PowerPoint (.pptx) files. "
        "Supports reading slide text, creating presentations, adding slides, "
        "replacing text, and listing presentation structure."
    )
    tool_parameters_schema = {
        "type": "object",
        "properties": {
            "action": {
                "type": "string",
                "enum": ["read", "create", "add_slide", "replace_text", "info"],
                "description": (
                    "Action to perform: "
                    "'read' extracts all text from slides, "
                    "'create' creates a new presentation, "
                    "'add_slide' adds a slide with title and content, "
                    "'replace_text' finds and replaces text across all slides, "
                    "'info' returns presentation metadata and structure."
                ),
            },
            "path": {
                "type": "string",
                "description": "Path to the .pptx file.",
            },
            "title": {
                "type": "string",
                "description": "Slide title (for create/add_slide).",
            },
            "content": {
                "type": "string",
                "description": "Slide body content (for create/add_slide). Use newlines for bullet points.",
            },
            "layout_index": {
                "type": "integer",
                "description": "Slide layout index (default: 1 = title and content).",
            },
            "find": {
                "type": "string",
                "description": "Text to find (for replace_text).",
            },
            "replace": {
                "type": "string",
                "description": "Replacement text (for replace_text).",
            },
        },
        "required": ["action", "path"],
    }

    tool_requires_approval = True
    tool_approval_risk_level = ApprovalRiskLevel.MEDIUM
    tool_supports_parallelism = True

    async def _execute(self, **kwargs: Any) -> dict[str, Any]:
        """Dispatch to the requested action handler."""
        action = kwargs["action"]
        handler = {
            "read": self._read,
            "create": self._create,
            "add_slide": self._add_slide,
            "replace_text": self._replace_text,
            "info": self._info,
        }.get(action)

        if handler is None:
            return {"success": False, "error": f"Unknown action: {action}"}

        return await handler(**kwargs)

    async def _read(self, **kwargs: Any) -> dict[str, Any]:
        """Read all text from a .pptx file."""
        pptx = _require_pptx()
        path = kwargs["path"]
        prs = pptx.Presentation(path)

        slides_data = []
        for i, slide in enumerate(prs.slides, 1):
            texts = []
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for paragraph in shape.text_frame.paragraphs:
                        text = paragraph.text.strip()
                        if text:
                            texts.append(text)
                if shape.has_table:
                    for row in shape.table.rows:
                        row_texts = [cell.text for cell in row.cells]
                        texts.append(" | ".join(row_texts))
            slides_data.append({"slide_number": i, "texts": texts})

        all_text = "\n\n".join(
            f"--- Slide {s['slide_number']} ---\n" + "\n".join(s["texts"])
            for s in slides_data
        )

        return {
            "success": True,
            "text": all_text,
            "slides": slides_data,
            "slide_count": len(slides_data),
        }

    async def _create(self, **kwargs: Any) -> dict[str, Any]:
        """Create a new .pptx presentation with an optional title slide."""
        pptx = _require_pptx()
        path = kwargs["path"]
        title = kwargs.get("title", "")
        content = kwargs.get("content", "")

        prs = pptx.Presentation()

        if title or content:
            layout = prs.slide_layouts[0]  # Title slide
            slide = prs.slides.add_slide(layout)
            if title and slide.placeholders[0]:
                slide.placeholders[0].text = title
            if content and len(slide.placeholders) > 1:
                slide.placeholders[1].text = content

        Path(path).parent.mkdir(parents=True, exist_ok=True)
        prs.save(path)
        logger.info("pptx.created", path=path)
        return {"success": True, "path": path, "message": f"Presentation created at {path}"}

    async def _add_slide(self, **kwargs: Any) -> dict[str, Any]:
        """Add a slide to an existing .pptx file."""
        pptx = _require_pptx()
        path = kwargs["path"]
        title = kwargs.get("title", "")
        content = kwargs.get("content", "")
        layout_index = kwargs.get("layout_index", 1)

        prs = pptx.Presentation(path)

        if layout_index >= len(prs.slide_layouts):
            layout_index = 1

        layout = prs.slide_layouts[layout_index]
        slide = prs.slides.add_slide(layout)

        if title and len(slide.placeholders) > 0:
            slide.placeholders[0].text = title
        if content and len(slide.placeholders) > 1:
            slide.placeholders[1].text = content

        prs.save(path)
        slide_num = len(prs.slides)
        logger.info("pptx.slide_added", path=path, slide_number=slide_num)
        return {
            "success": True,
            "path": path,
            "slide_number": slide_num,
            "message": f"Slide {slide_num} added",
        }

    async def _replace_text(self, **kwargs: Any) -> dict[str, Any]:
        """Find and replace text across all slides."""
        pptx = _require_pptx()
        path = kwargs["path"]
        find_text = kwargs.get("find", "")
        replace_text = kwargs.get("replace", "")

        if not find_text:
            return {"success": False, "error": "Parameter 'find' is required for replace_text"}

        prs = pptx.Presentation(path)
        count = 0

        for slide in prs.slides:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for paragraph in shape.text_frame.paragraphs:
                        for run in paragraph.runs:
                            if find_text in run.text:
                                run.text = run.text.replace(find_text, replace_text)
                                count += 1

        prs.save(path)
        logger.info("pptx.text_replaced", path=path, replacements=count)
        return {
            "success": True,
            "path": path,
            "replacements": count,
            "message": f"Replaced {count} occurrence(s)",
        }

    async def _info(self, **kwargs: Any) -> dict[str, Any]:
        """Return presentation metadata and structure."""
        pptx = _require_pptx()
        path = kwargs["path"]
        prs = pptx.Presentation(path)

        slides_info = []
        for i, slide in enumerate(prs.slides, 1):
            shapes = []
            for shape in slide.shapes:
                shape_info: dict[str, Any] = {
                    "name": shape.name,
                    "shape_type": str(shape.shape_type),
                }
                if shape.has_text_frame:
                    shape_info["has_text"] = True
                    shape_info["text_preview"] = shape.text_frame.text[:100]
                if shape.has_table:
                    shape_info["has_table"] = True
                    shape_info["table_rows"] = len(shape.table.rows)
                    shape_info["table_cols"] = len(shape.table.columns)
                shapes.append(shape_info)
            slides_info.append({"slide_number": i, "shapes": shapes})

        core_props = prs.core_properties
        return {
            "success": True,
            "slide_count": len(prs.slides),
            "layout_count": len(prs.slide_layouts),
            "slides": slides_info,
            "author": core_props.author or "",
            "title": core_props.title or "",
            "created": str(core_props.created) if core_props.created else "",
            "modified": str(core_props.modified) if core_props.modified else "",
        }
